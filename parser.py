from pydantic import BaseModel
from typing import List
from pptx import Presentation
import os
from dotenv import load_dotenv
from typing import Dict
import json 
from openai import OpenAI
import ast
import pymupdf
import base64
from pdf2image import convert_from_path
import subprocess

# JSON Object Definitions
class DetailsTable(BaseModel):
    observation: str
    risk: str
    risk_rating: str
    recommendation: str
    status: str

class ReportData(BaseModel):
    report_title : str
    executive_summary : str
    overall_audit_rating : str 
    overall_risk_description : str
    details: List[DetailsTable]
    recommendations: List[str]
    management_action_plan: List[str]

class BatchResponse(BaseModel):
    reports: List[ReportData]
    
def extract_slide_content(file_path):
    prs = Presentation(file_path)
    slides_content = []

    for i, slide in enumerate(prs.slides):
        slide_text_blocks = []
        
        for shape in slide.shapes:
            # Extract from standard text shapes
            if hasattr(shape, "text") and shape.text.strip():
                slide_text_blocks.append(shape.text.strip())
            
            # Extract from tables
            elif shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    # Join cell text with a pipe (|) to maintain table structure
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        table_data.append(row_text)
                
                if table_data:
                    slide_text_blocks.append("\n".join(table_data))

        # Join all text found on this specific slide into one string
        full_slide_text = "\n".join(slide_text_blocks)
        slides_content.append(f"[SLIDE START]\n{full_slide_text}\n[SLIDE END]\n")
    
    return "\n".join(slides_content)

def extract_pdf_content(file_path):
    doc = pymupdf.open(file_path)
    full_text = ""
    for page in doc:
        full_text += f"[SLIDE START]\n{page.get_text("text")}\n[SLIDE_END]\n"
    return full_text

def generate_json(extracted_text):
    # Load extraction prompt
    extraction_prompt = ""
    try:
        with open('prompts/parser_prompt.txt', 'r') as f:
            extraction_prompt = f.read()
    except FileNotFoundError:
        print("ERROR: System prompt text file not found...")
        exit(-1)
    
    # Load and import API Keys
    load_dotenv()
    OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
    
    openai_client = OpenAI(
        api_key=OPENAI_API_KEY
    )
    
    # Parse text content from PPT
    response = openai_client.responses.create(
        model='gpt-5-mini', 
        instructions=extraction_prompt,
        input=extracted_text,
    )
    
    raw = response.output_text.strip()
    data = ast.literal_eval(raw)
    print(data)
    return data


def get_overall_rating_slide_image(pptx_path, output_path):
    print("Converting PPTX to PDF via LibreOffice...")
    # Silently command LibreOffice to convert the file
    subprocess.run([
        "libreoffice", "--headless", "--convert-to", "pdf",
        "--outdir", os.path.dirname(os.path.abspath(pptx_path)),
        pptx_path
    ], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    
    # Get the path of the new PDF
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    pdf_path = os.path.join(os.path.dirname(pptx_path), f"{base_name}.pdf")
    slide_number = 3
    print(f"Extracting Slide {slide_number}...")
    # Extract the specific page (1-based index)
    images = convert_from_path(pdf_path, first_page=slide_number, last_page=slide_number)
    
    if images:
        images[0].save(output_path, 'PNG')
        print(f"Success! Image saved to {output_path}")
    
    # Clean up the PDF
    if os.path.exists(pdf_path):
        os.remove(pdf_path)
        
        
def get_overall_rating_pdf_image(pdf_path, output_path):
    doc = pymupdf.open(pdf_path)
    page = doc.load_page(2)
    pix = page.get_pixmap(dpi=150)
    pix.save(output_path)
    doc.close()
    
def base64_encode_image(image_path):
    with open(image_path, 'rb') as image:
        return base64.b64encode(image.read()).decode('utf-8')
    
def get_overall_rating_from_image(image_path):
    # Encode image to base64
    base64_image = base64_encode_image(image_path)
    
    # Get image analysis prompt
    analysis_prompt = ""
    try:
        with open('prompts/image_analysis_prompt.txt', 'r') as f:
            analysis_prompt = f.read()
    except FileNotFoundError:
        print("ERROR: System prompt text file not found...")
        exit(-1)

    # Extract highlighted text using OpenAI API
    load_dotenv()
    OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
    
    openai_client = OpenAI(
        api_key=OPENAI_API_KEY
    )
    
    # Parse text content from PPT
    response = openai_client.responses.create(
        model='gpt-5-mini', 
        instructions=analysis_prompt,
        input=[
            {
                "role": "user",
                "content": [
                    {"type": "input_text", "text": "Find the highlighted text in the image."},
                    {
                        "type": "input_image",
                        "image_url": f"data:image/jpeg;base64,{base64_image}"
                    }
                ]
            }
            ],
    )
    print(f"Overall Rating: {response.output_text}")
    return response.output_text
    
def main():
    batch_json = []
    all_pptx = [f for f in os.listdir('generated_reports') if f.endswith('.pptx') or f.endswith('.pdf')]
    for i in range(0, len(all_pptx), 5):
        combined_text = ""
        overall_rating = ""
        batch_files = all_pptx[i:i+5]
        for filename in batch_files:
            text = ""
            if filename.endswith(".pdf"):
                text = extract_pdf_content(f"generated_reports/{filename}")
                # get_overall_rating_pdf_image(f"generated_reports/{filename}", f"images/{filename}.jpg")
            else:
                text = extract_slide_content(f"generated_reports/{filename}")
                # get_overall_rating_slide_image(f"generated_reports/{filename}", f"images/{filename}.jpg")
            # overall_rating = get_overall_rating_from_image(f'images/{filename}.jpg')
            combined_text += f"\n[[START_FILE: {filename}]]\n{text}\n[[END_FILE]]\n"
            print(combined_text)
        # Create Batch JSONs and add it to final JSON list
        batch_json += generate_json(combined_text)

    # Create JSON file
    with open("output_json/parsed_output.json", "w") as f:
        json.dump(batch_json, f, indent=4)
if __name__ == '__main__':
    main()