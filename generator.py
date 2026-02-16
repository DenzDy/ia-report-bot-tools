from openai import OpenAI
import os
from dotenv import load_dotenv
import sys
import argparse
import json
import markdown2
import subprocess

from pptx import Presentation
import os

def fill_slide_by_layout_names(slide, data_dict):
    """
    Maps layout placeholder names to the slide placeholder indices.
    """
    # Create a map of {Name: Index} from the LAYOUT
    layout_map = {ph.name: ph.placeholder_format.idx for ph in slide.slide_layout.placeholders}
    
    # Use that map to fill the placeholders on the SLIDE
    for name, idx in layout_map.items():
        if name in data_dict:
            # Access the placeholder on the slide using the index found on the layout
            slide.placeholders[idx].text = str(data_dict[name])
            print(f"Filled {name} (Index {idx})")
            
def export_as_pptx(json_data, target_directory="generated_reports"):
    if not os.path.exists(target_directory):
        os.makedirs(target_directory)

    for report in json_data:
        # Load your template
        # Ensure 'template.pptx' is in your directory
        prs = Presentation('templates/ac_template.pptx') 
        
        report_fn = report['file_name'].replace('.json', '.pptx')
        output_path = os.path.join(target_directory, report_fn)
        
        # Placeholder Filling for Slides
        
        # Report Title Slide
        report_content = report['Report Content']
        current_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(current_layout)
        fill_slide_by_layout_names(slide, report_content['Report Title Slide'])
        
        # Scope, Background Slide
        current_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(current_layout)
        fill_slide_by_layout_names(slide, report_content['Objective, Background and Scope'])
        
        # Overall Audit Rating Slide
        current_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(current_layout)
        fill_slide_by_layout_names(slide, report_content['Overall Internal Audit Rating'])
        
        # Observation Slides
        current_layout = prs.slide_layouts[3]
        for obs in report_content['Observation Slides']:
            slide = prs.slides.add_slide(current_layout)
            fill_slide_by_layout_names(slide, obs)
        
        # End of Report Slide 
        current_layout = prs.slide_layouts[4]
        slide = prs.slides.add_slide(current_layout)
        
        # Save the file
        prs.save(output_path)
        print(f"Success! Created {output_path} using python-pptx")
def main():
    # Load and import API Keys
    load_dotenv()
    OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')

    # Parse input arguments
    parser = argparse.ArgumentParser(description="Generates a JSON of internal audit risk reports in Markdown format.")
    parser.add_argument("-r","--reports", required=True, type=int, help="number of reports to generate.")
    parser.add_argument("-v", '--verbose', action='store_true')
    parser.add_argument("--seed",
        type=str,
        help="Optional: Provide a specific industry or focus (e.g., 'Healthcare', 'Fintech')",
        default="General Corporate Operations"
    )
    parser.add_argument("--type",
        type=str,
        help="Optional: Provide a specific report type",
        default="General Corporate Operations"
    )
    parser.add_argument("--observations",
        type=int,
        help="Optional: Provide the number of observations per report",
        default=2
    )
    args = parser.parse_args()

    # Verbose logging of command line input arguments
    if args.verbose:
        print(f"[DEBUG] Received reports={args.reports}")

    # Get response from Gemini API
    openai_client = OpenAI(
        api_key=OPENAI_API_KEY
    )

    # Load System Prompt
    system_prompt = ""
    try:
        with open('prompts/generation_prompt.txt', 'r') as f:
            system_prompt = f.read()
    except FileNotFoundError:
        print("ERROR: System prompt text file not found...")
        exit(-1)
        
    # User Prompt
    user_prompt = f"""
    Generate {args.reports} reports with focus on {args.seed} for the reports with {args.observations} observations for each report.
    """
    response = openai_client.responses.create(
        model='gpt-5-mini', 
        instructions=system_prompt,
        input=user_prompt
    )
    # print(response.output_text)
    # Load JSON
    data = json.loads(response.output_text)
    print(data)
    with open("output.json", "w", encoding="utf-8") as json_file:
        json.dump(data, json_file, ensure_ascii=False, indent=4)
    # Convert data to PDFs
    # DEBUG: Load dummy response from JSON file
    with open('output.json', 'r') as file:
        data = json.load(file)
    export_as_pptx(data)
    
if __name__ == '__main__':
    main()
